"""
Premium PowerPoint Builder Module
Creates stunning presentations with custom shapes, accent bars, and themed styling.
Inspired by the NotebookLM presentation style.

All slides use blank layouts with custom-positioned shapes for maximum visual control.
Minimum text size: 16pt for readability during presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from colorama import Fore, Style
from modules.themes import get_theme, PREMIUM_THEMES


# ============================================================================
# SLIDE DIMENSIONS (16:9 Widescreen)
# ============================================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============================================================================
# LOW-LEVEL SHAPE HELPERS
# ============================================================================

def set_slide_bg(slide, color):
    """Set a solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, corner_radius=0.05):
    """Add a rounded rectangle (card) with fill color."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # No border
    if shape.adjustments:
        shape.adjustments[0] = corner_radius
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    """Add a thin colored accent bar (decorative line)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(255,255,255), bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=RGBColor(180,185,200), bullet_color=RGBColor(0,206,209),
                    font_name="Segoe UI", line_spacing=1.4):
    """Add a bullet list with colored chevron markers and generous spacing."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Colored bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "  >  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        run_bullet.font.name = font_name

        # Content text
        run_text = p.add_run()
        run_text.text = str(item)
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name

        p.space_after = Pt(6)
        p.space_before = Pt(3)
        p.line_spacing = line_spacing

    return txBox


# ============================================================================
# PREMIUM SLIDE BUILDERS
# ============================================================================

def add_title_slide(prs, topic, theme):
    """Create a stunning title slide with accent bars and styled text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, theme['bg_primary'])

    # Top accent bar
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), theme['accent_1'])

    # Subtitle tag
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.4),
                 "AI-POWERED PRESENTATION", font_size=16,
                 color=theme['accent_2'], bold=True, font_name=theme['font_heading'])

    # Main title
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
                 topic, font_size=44, color=theme['text_primary'], bold=True,
                 font_name=theme['font_heading'])

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
                 "Comprehensive Research-Based Presentation\nGenerated from AI Analysis",
                 font_size=20, color=theme['text_secondary'], font_name=theme['font_body'])

    # Divider accent
    add_accent_bar(slide, Inches(1.5), Inches(5.3), Inches(2.5), Inches(0.04), theme['accent_3'])

    # Author line
    add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10), Inches(0.4),
                 "Generated with FlashPoint AI  |  Python + Gemini",
                 font_size=16, color=theme['text_dim'], font_name=theme['font_body'])

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), theme['accent_3'])


def add_toc_slide(prs, section_titles, theme):
    """Create a Table of Contents slide with numbered items in cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "Table of Contents", font_size=36, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), theme['title_bar_color'])

    # TOC items as cards
    accent_colors = [theme['accent_1'], theme['accent_2'], theme['accent_3']]

    for i, title in enumerate(section_titles[:10]):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(1.5) + row * Inches(1.1)

        color = accent_colors[i % len(accent_colors)]

        card = add_shape_rect(slide, x, y, Inches(5.8), Inches(0.85), theme['bg_card'])
        add_accent_bar(slide, x, y, Inches(0.06), Inches(0.85), color)

        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.color.rgb = theme['text_primary']
        p.font.bold = True
        p.font.name = theme['font_heading']
        p.alignment = PP_ALIGN.CENTER

        # Title text
        add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(4.5), Inches(0.5),
                     str(title), font_size=18, color=theme['text_primary'],
                     bold=False, font_name=theme['font_body'])


def add_content_slide(prs, title, bullets, notes="", theme=None, max_bullets_per_slide=6, explanation=""):
    """
    Create premium content slide(s) with accent bar, card background, styled bullets,
    and optional explanation text below bullets.
    Optimized for 6 bullets of up to 20 words each.
    """
    if theme is None:
        theme = PREMIUM_THEMES['1']  # Default to Midnight

    if not bullets:
        bullets = ["No content available"]

    # Split into chunks
    bullet_chunks = []
    current_chunk = []

    for bullet in bullets:
        bullet_text = str(bullet).strip()
        if len(bullet_text) > 300:
            bullet_text = bullet_text[:297] + "..."
        current_chunk.append(bullet_text)
        if len(current_chunk) >= max_bullets_per_slide:
            bullet_chunks.append(current_chunk)
            current_chunk = []

    if current_chunk:
        bullet_chunks.append(current_chunk)

    slides_created = []
    accent_colors = [theme['accent_1'], theme['accent_2'], theme['accent_3']]

    for chunk_idx, bullet_chunk in enumerate(bullet_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        set_slide_bg(slide, theme['bg_primary'])

        # Slide title
        slide_title = title if chunk_idx == 0 else f"{title} (continued)"

        add_text_box(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
                     slide_title, font_size=28, color=theme['text_primary'],
                     bold=True, font_name=theme['font_heading'])

        # Accent bar under title
        add_accent_bar(slide, Inches(0.8), Inches(0.95), Inches(3), Inches(0.04),
                       theme['title_bar_color'])

        # Content card background
        card_top = Inches(1.15)
        card_height = Inches(6.0)
        add_shape_rect(slide, Inches(0.6), card_top, Inches(12.1), card_height, theme['bg_card'])

        # Adjust bullet box height based on whether explanation exists
        has_explanation = bool(explanation) and chunk_idx == 0
        bullet_height = Inches(4.2) if has_explanation else Inches(5.7)

        # Bullet list
        bullet_color = accent_colors[chunk_idx % len(accent_colors)]
        add_bullet_list(slide, Inches(0.9), Inches(1.3), Inches(11.5), bullet_height,
                        bullet_chunk, font_size=16, color=theme['text_secondary'],
                        bullet_color=bullet_color, font_name=theme['font_body'],
                        line_spacing=1.3)

        # Explanation text below bullets (only on first chunk)
        if has_explanation:
            explanation_text = str(explanation).strip()
            if len(explanation_text) > 500:
                explanation_text = explanation_text[:497] + "..."
            
            # Divider line before explanation
            explanation_top = Inches(5.6)
            add_accent_bar(slide, Inches(1.2), explanation_top, Inches(10.9), Inches(0.02),
                          theme['accent_3'])
            
            # Explanation text — italic, slightly smaller, dimmer color
            txBox = slide.shapes.add_textbox(
                Inches(1.2), Inches(5.75), Inches(10.9), Inches(1.3)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = explanation_text
            run.font.size = Pt(13)
            run.font.italic = True
            run.font.color.rgb = theme.get('text_dim', theme['text_secondary'])
            run.font.name = theme['font_body']
            p.line_spacing = 1.2

        # Speaker notes — disabled per user request
        # if notes and chunk_idx == 0:
        #     notes_slide = slide.notes_slide
        #     text_frame = notes_slide.notes_text_frame
        #     text_frame.text = str(notes)

        slides_created.append(slide)

    return slides_created


def add_image_slide(prs, title, image_path, caption="", theme=None):
    """Add a themed slide with an image."""
    if theme is None:
        theme = PREMIUM_THEMES['1']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=32, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04),
                   theme['title_bar_color'])

    # Add image (centered in card)
    try:
        card = add_shape_rect(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(5.5), theme['bg_card'])

        from PIL import Image
        img = Image.open(image_path)
        img_width, img_height = img.size
        aspect = img_width / img_height

        # Fit within card
        max_w = Inches(10)
        max_h = Inches(5)

        if aspect > 2:
            width = max_w
            height = int(width / aspect)
        else:
            height = max_h
            width = int(height * aspect)

        # Center image
        left = Inches(1.5) + (max_w - width) // 2
        top = Inches(1.5) + (max_h - height) // 2

        slide.shapes.add_picture(image_path, left, top, width, height)

    except Exception as e:
        add_text_box(slide, Inches(1.5), Inches(3), Inches(10), Inches(1),
                     f"[Image: {os.path.basename(image_path)}]",
                     font_size=18, color=theme['text_dim'], font_name=theme['font_body'])

    # Caption
    if caption:
        add_text_box(slide, Inches(1), Inches(6.9), Inches(11.3), Inches(0.4),
                     caption, font_size=16, color=theme['text_dim'],
                     alignment=PP_ALIGN.CENTER, font_name=theme['font_body'])


def add_two_column_slide(prs, title, left_bullets, right_bullets, theme=None):
    """Create a two-column content slide with card backgrounds."""
    if theme is None:
        theme = PREMIUM_THEMES['1']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=32, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04),
                   theme['title_bar_color'])

    # Left card
    add_shape_rect(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.8), theme['bg_card'])
    add_accent_bar(slide, Inches(0.5), Inches(1.3), Inches(0.06), Inches(5.8), theme['accent_1'])
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(5.3), Inches(5.2),
                    left_bullets, font_size=16, color=theme['text_secondary'],
                    bullet_color=theme['accent_1'], font_name=theme['font_body'])

    # Right card
    add_shape_rect(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(5.8), theme['bg_card'])
    add_accent_bar(slide, Inches(6.8), Inches(1.3), Inches(0.06), Inches(5.8), theme['accent_2'])
    add_bullet_list(slide, Inches(7.1), Inches(1.5), Inches(5.3), Inches(5.2),
                    right_bullets, font_size=16, color=theme['text_secondary'],
                    bullet_color=theme['accent_2'], font_name=theme['font_body'])


def add_key_takeaways_slide(prs, key_points, theme=None):
    """Create a summary/key takeaways slide with individual cards per point."""
    if theme is None:
        theme = PREMIUM_THEMES['1']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "Key Takeaways", font_size=36, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04),
                   theme['title_bar_color'])

    accent_colors = [theme['accent_1'], theme['accent_2'], theme['accent_3']]

    for i, point in enumerate(key_points[:5]):
        y = Inches(1.4) + i * Inches(1.15)
        color = accent_colors[i % len(accent_colors)]

        card = add_shape_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.95), theme['bg_card'])
        add_accent_bar(slide, Inches(0.8), y, Inches(0.06), Inches(0.95), color)

        # Number
        add_text_box(slide, Inches(1.1), y + Inches(0.2), Inches(0.5), Inches(0.5),
                     str(i + 1), font_size=22, color=color, bold=True,
                     font_name=theme['font_heading'])

        # Point text
        point_text = str(point)
        if len(point_text) > 150:
            point_text = point_text[:147] + "..."
        add_text_box(slide, Inches(1.8), y + Inches(0.2), Inches(10.3), Inches(0.6),
                     point_text, font_size=18, color=theme['text_secondary'],
                     font_name=theme['font_body'])


def add_closing_slide(prs, topic, theme=None):
    """Create a professional closing/thank you slide."""
    if theme is None:
        theme = PREMIUM_THEMES['1']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    # Top accent
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), theme['accent_3'])

    # Tag
    add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(0.4),
                 "AI-POWERED PRESENTATION", font_size=16, color=theme['accent_2'],
                 bold=True, font_name=theme['font_heading'])

    # Thank You
    add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(1),
                 "Thank You!", font_size=52, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
                 f"This presentation on '{topic}' was generated using\nAI-powered research and professional formatting.",
                 font_size=20, color=theme['text_secondary'], font_name=theme['font_body'])

    add_accent_bar(slide, Inches(1.5), Inches(5.2), Inches(2), Inches(0.04), theme['accent_1'])

    # Footer
    import datetime
    details = [
        "Generated with FlashPoint AI",
        f"Research Date: {datetime.datetime.now().strftime('%B %Y')}",
        "Questions & Discussion Welcome",
    ]
    add_bullet_list(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(1.5),
                    details, font_size=16, color=theme['text_dim'],
                    bullet_color=theme['accent_3'], font_name=theme['font_body'])

    # Bottom accent
    add_accent_bar(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), theme['accent_1'])


def add_references_slide(prs, theme=None):
    """Create a References slide."""
    if theme is None:
        theme = PREMIUM_THEMES['1']

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme['bg_primary'])

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 "References", font_size=36, color=theme['text_primary'],
                 bold=True, font_name=theme['font_heading'])
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04),
                   theme['title_bar_color'])

    import datetime
    refs = [
        "Research compiled using Google Gemini AI",
        "Web sources and academic publications",
        f"Current as of {datetime.datetime.now().strftime('%B %Y')}",
        "All information sourced from reputable sources",
    ]

    add_shape_rect(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), theme['bg_card'])
    add_bullet_list(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5),
                    refs, font_size=18, color=theme['text_secondary'],
                    bullet_color=theme['accent_2'], font_name=theme['font_body'],
                    line_spacing=1.8)


# ============================================================================
# MAIN BUILD FUNCTION
# ============================================================================

def build_presentation(topic, slides, image_paths, output_filename, config=None):
    """
    Build a premium PowerPoint presentation.

    Uses the theme from config to apply consistent styling across all slides.
    All slides use blank layouts with custom shapes for maximum visual control.
    """
    print(f"{Fore.CYAN}[*] Building premium PowerPoint presentation...{Style.RESET_ALL}")

    # Get theme
    theme_choice = '1'  # Default: Midnight
    if config and config.get('theme_choice'):
        theme_choice = str(config['theme_choice'])
    theme = get_theme(theme_choice)
    print(f"{Fore.CYAN}  Theme: {theme['name']}{Style.RESET_ALL}")

    # Create presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_count = 0
    image_index = 0
    
    # Build debug log
    debug_lines = []
    debug_lines.append(f"=== PPT BUILD DEBUG LOG ===")
    debug_lines.append(f"Topic: {topic}")
    debug_lines.append(f"Theme: {theme['name']}")
    debug_lines.append(f"Total slides received: {len(slides)}")
    debug_lines.append(f"Images available: {len(image_paths)}")
    debug_lines.append(f"{'='*70}\n")

    # Collect section titles for TOC
    toc_titles = []
    for s in slides:
        title = s.get('title', '')
        if title and title not in ('Table of Contents', topic, 'Key Takeaways',
                                   'Conclusion', 'References') and '(continued)' not in title:
            if title not in toc_titles:
                toc_titles.append(title)

    for i, slide_data in enumerate(slides):
        title = slide_data.get('title', 'Untitled')
        bullets = slide_data.get('bullets', [])
        notes = slide_data.get('notes', '')
        explanation = slide_data.get('explanation', '')
        slide_type = slide_data.get('type', 'CONTENT')

        # Debug log entry
        debug_lines.append(f"--- BUILD SLIDE {i+1} ---")
        debug_lines.append(f"  Type:        {slide_type}")
        debug_lines.append(f"  Title:       {title}")
        debug_lines.append(f"  Bullets ({len(bullets)}):")
        for j, b in enumerate(bullets):
            debug_lines.append(f"    [{j+1}] {str(b)[:120]}")
        debug_lines.append(f"  Explanation: {explanation[:120]}{'...' if len(explanation) > 120 else ''}")
        debug_lines.append(f"  Notes:       {'YES' if notes else 'NO'} ({len(notes)} chars)")

        # Skip empty slides
        if not bullets and not notes and i > 1:
            debug_lines.append(f"  ACTION:      *** SKIPPED (no bullets/notes) ***\n")
            continue

        # Title slide
        if i == 0:
            add_title_slide(prs, topic, theme)
            slide_count += 1
            debug_lines.append(f"  ACTION:      Title slide created\n")
            continue

        # Table of Contents
        if title == 'Table of Contents':
            add_toc_slide(prs, toc_titles[:10], theme)
            slide_count += 1
            debug_lines.append(f"  ACTION:      TOC slide created\n")
            continue

        # Key Takeaways
        if title == 'Key Takeaways':
            add_key_takeaways_slide(prs, bullets, theme)
            slide_count += 1
            debug_lines.append(f"  ACTION:      Key Takeaways slide created\n")
            continue

        # Conclusion / Thank You
        if title in ('Conclusion', 'Thank You'):
            add_closing_slide(prs, topic, theme)
            slide_count += 1
            debug_lines.append(f"  ACTION:      Closing slide created\n")
            continue

        # References
        if title == 'References':
            add_references_slide(prs, theme)
            slide_count += 1
            debug_lines.append(f"  ACTION:      References slide created\n")
            continue

        # Regular content slide
        created = add_content_slide(prs, title, bullets, notes, theme, explanation=explanation)
        if isinstance(created, list):
            slide_count += len(created)
            debug_lines.append(f"  ACTION:      Content slide created ({len(created)} actual slides)")
        else:
            slide_count += 1
            debug_lines.append(f"  ACTION:      Content slide created (1 actual slide)")
        
        debug_lines.append(f"  RENDERED:    {len(bullets)} bullets, explanation={'YES' if explanation else 'NO'}, notes={'YES' if notes else 'NO'}\n")

        # Insert image slides periodically
        if image_index < len(image_paths) and (i % 4 == 0 or i % 3 == 0) and i > 2:
            image_title = f"{topic} - Visual {image_index + 1}"
            add_image_slide(prs, image_title, image_paths[image_index], theme=theme)
            slide_count += 1
            image_index += 1
            print(f"{Fore.GREEN}  [+] Added image slide: {os.path.basename(image_paths[image_index - 1])}{Style.RESET_ALL}")

    # Add remaining images
    while image_index < len(image_paths):
        image_title = f"{topic} - Visual {image_index + 1}"
        add_image_slide(prs, image_title, image_paths[image_index], theme=theme)
        slide_count += 1
        image_index += 1
        print(f"{Fore.GREEN}  [+] Added image slide: {os.path.basename(image_paths[image_index - 1])}{Style.RESET_ALL}")

    print(f"{Fore.GREEN}[+] Premium presentation built!{Style.RESET_ALL}")
    print(f"{Fore.CYAN}  Total slides: {slide_count}{Style.RESET_ALL}")
    print(f"{Fore.CYAN}  Theme: {theme['name']}{Style.RESET_ALL}")
    print(f"{Fore.CYAN}  Images included: {image_index}{Style.RESET_ALL}")
    
    # Save build debug log
    debug_lines.append(f"\n{'='*70}")
    debug_lines.append(f"TOTAL PPT SLIDES BUILT: {slide_count}")
    debug_lines.append(f"IMAGES INCLUDED: {image_index}")
    try:
        output_dir = os.path.dirname(output_filename)
        if output_dir:
            log_path = os.path.join(output_dir, 'ppt_build_debug.txt')
            with open(log_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(debug_lines))
            print(f"{Fore.GREEN}  [+] Build debug saved to: {log_path}{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.YELLOW}  [!] Could not save build debug log: {e}{Style.RESET_ALL}")

    return prs
