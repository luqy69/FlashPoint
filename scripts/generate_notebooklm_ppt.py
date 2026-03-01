"""
Generate a NotebookLM-style presentation explaining how the AI PPT Generator works.
Creates a professional, dark-themed presentation with modern styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette (NotebookLM-inspired dark theme) ─────────────────────────
BG_DARK       = RGBColor(15, 15, 25)       # Deep space blue-black
BG_CARD       = RGBColor(30, 32, 48)       # Card backgrounds
ACCENT_BLUE   = RGBColor(100, 149, 237)    # Cornflower blue
ACCENT_PURPLE = RGBColor(147, 112, 219)    # Medium purple
ACCENT_CYAN   = RGBColor(0, 206, 209)      # Dark turquoise
ACCENT_GREEN  = RGBColor(0, 200, 150)      # Emerald
ACCENT_ORANGE = RGBColor(255, 165, 0)      # Gold
TEXT_WHITE     = RGBColor(240, 240, 248)    # Near-white
TEXT_GRAY      = RGBColor(180, 185, 200)    # Soft gray
TEXT_DIM       = RGBColor(120, 125, 145)    # Dimmed text
GRADIENT_START = RGBColor(25, 25, 60)       # Gradient purple-ish
GRADIENT_END   = RGBColor(10, 10, 20)       # Almost black

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set a solid dark background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    """Add a rounded rectangle with fill color."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # No border
    # Set corner radius via adjustment
    if shape.adjustments:
        shape.adjustments[0] = 0.05  # subtle rounding
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    """Add a thin colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_GRAY, bullet_color=ACCENT_CYAN, font_name="Segoe UI"):
    """Add a bullet list with colored bullets."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "  >  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        run_bullet.font.name = font_name

        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name

        p.space_after = Pt(8)

    return txBox


def add_numbered_step(slide, left, top, number, title, description,
                      accent_color=ACCENT_BLUE, width=Inches(11)):
    """Add a numbered step block (number badge + title + description)."""
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent_color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    # Title
    add_text_box(slide, left + Inches(0.75), top, width - Inches(0.75), Inches(0.4),
                 title, font_size=20, color=TEXT_WHITE, bold=True)

    # Description
    add_text_box(slide, left + Inches(0.75), top + Inches(0.4), width - Inches(0.75), Inches(0.5),
                 description, font_size=14, color=TEXT_GRAY)


# ════════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)

    # Top accent line
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    # Subtitle / tag
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.4),
                 "NOTEBOOKLM  |  SOURCE GUIDE", font_size=14, color=ACCENT_CYAN, bold=True)

    # Main title
    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
                 "AI PowerPoint Presentation Generator", font_size=44,
                 color=TEXT_WHITE, bold=True)

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
                 "How it automates research, content creation,\nand professional slide design using Google Gemini",
                 font_size=20, color=TEXT_GRAY)

    # Divider
    add_accent_bar(slide, Inches(1.5), Inches(4.8), Inches(2), Inches(0.04), ACCENT_PURPLE)

    # Author / meta
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
                 "Made by Luqman  (2024-ag-8738 DVM)   |   Python + Selenium + Gemini AI",
                 font_size=14, color=TEXT_DIM)

    # Bottom gradient bar
    add_accent_bar(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), ACCENT_PURPLE)


def slide_overview(prs):
    """Slide 2: Program overview / What it does."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "What is FlashPoint AI?", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_CYAN)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(1),
                 "An intelligent Python application that automatically generates professional\n"
                 "PowerPoint presentations on ANY topic using Google Gemini's deep research\n"
                 "capabilities and browser automation.",
                 font_size=18, color=TEXT_GRAY)

    # Feature cards (2-column grid)
    features = [
        ("Automated Research", "Uses Selenium to control Chrome and\nleverage Gemini's deep research mode", ACCENT_BLUE),
        ("Professional Design", "Clean, modern slide layouts with\nconsistent formatting and themes", ACCENT_PURPLE),
        ("Data Visualizations", "Generates charts, graphs and downloads\nrelevant images from Google", ACCENT_GREEN),
        ("Speaker Notes", "Includes detailed speaker notes\nfor each slide automatically", ACCENT_ORANGE),
        ("Super Research", "Integrates academic APIs: CORE,\narXiv, PubMed, OpenLibrary", ACCENT_CYAN),
        ("Multi-Format Export", "Export to PPTX, PPT (97-2003),\nand PDF formats", ACCENT_BLUE),
    ]

    for i, (title, desc, color) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6)
        y = Inches(2.8) + row * Inches(1.5)

        card = add_shape_rect(slide, x, y, Inches(5.5), Inches(1.3), BG_CARD)
        add_accent_bar(slide, x, y, Inches(0.06), Inches(1.3), color)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5), Inches(0.35),
                     title, font_size=18, color=color, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(5), Inches(0.7),
                     desc, font_size=13, color=TEXT_GRAY)


def slide_architecture(prs):
    """Slide 3: System architecture overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "System Architecture", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_PURPLE)

    # Entry points
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.35),
                 "ENTRY POINTS", font_size=14, color=ACCENT_CYAN, bold=True)

    entries = [
        ("ppt_generator.py", "CLI interface  -  Terminal-based workflow", ACCENT_BLUE),
        ("ppt_wizard.py", "GUI interface  -  Tkinter wizard (6 steps)", ACCENT_PURPLE),
    ]
    for i, (name, desc, color) in enumerate(entries):
        y = Inches(1.9) + i * Inches(0.8)
        card = add_shape_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.65), BG_CARD)
        add_accent_bar(slide, Inches(0.8), y, Inches(0.05), Inches(0.65), color)
        add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.3),
                     name, font_size=15, color=color, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(1.1), y + Inches(0.35), Inches(5), Inches(0.3),
                     desc, font_size=12, color=TEXT_GRAY)

    # Core modules
    add_text_box(slide, Inches(0.8), Inches(3.7), Inches(5), Inches(0.35),
                 "CORE MODULES  (modules/)", font_size=14, color=ACCENT_CYAN, bold=True)

    modules = [
        ("research.py", "Gemini browser automation (928 lines)", ACCENT_BLUE),
        ("content_generator.py", "Outline & slide structuring (333 lines)", ACCENT_GREEN),
        ("ppt_builder.py", "PowerPoint assembly (373 lines)", ACCENT_PURPLE),
        ("visualizations.py", "Image download & chart gen (284 lines)", ACCENT_ORANGE),
        ("themes.py", "10 built-in color themes (219 lines)", ACCENT_CYAN),
        ("export.py", "PPT & PDF export via COM (143 lines)", ACCENT_BLUE),
        ("research_doc.py", "Supplemental Word doc (374 lines)", ACCENT_GREEN),
        ("scientific_research.py", "Academic APIs (559 lines)", ACCENT_PURPLE),
    ]

    for i, (name, desc, color) in enumerate(modules):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.2)
        y = Inches(4.1) + row * Inches(0.7)
        card = add_shape_rect(slide, x, y, Inches(5.8), Inches(0.55), BG_CARD)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.8), Inches(0.25),
                     name, font_size=13, color=color, bold=True, font_name="Consolas")
        add_text_box(slide, x + Inches(2.8), y + Inches(0.05), Inches(2.8), Inches(0.25),
                     desc, font_size=11, color=TEXT_DIM)


def slide_pipeline(prs):
    """Slide 4: 5-Step pipeline flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
                 "How It Works  -  The 5-Step Pipeline", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3), Inches(0.04), ACCENT_GREEN)

    steps = [
        (1, "Web Research via Gemini", 
         "Selenium opens Chrome, navigates to gemini.google.com, submits a research prompt, "
         "waits for deep research to complete, then extracts structured content.",
         ACCENT_BLUE),
        (2, "Content Generation",
         "ContentGenerator analyzes research data, creates an outline with sections, "
         "bullet points, and speaker notes for 10-30 slides.",
         ACCENT_PURPLE),
        (3, "Visualization Creation",
         "Downloads relevant images from Google Images via Selenium. Falls back to "
         "matplotlib sample charts (bar, pie, line) if download fails.",
         ACCENT_GREEN),
        (4, "PowerPoint Assembly",
         "python-pptx builds the .pptx file: title slide, table of contents, content "
         "slides, image slides, references, and thank-you slide.",
         ACCENT_ORANGE),
        (5, "Theme & Export",
         "Applies one of 10 built-in themes (Ion, Integral, Facet, etc.) and "
         "optionally exports to PPT or PDF via PowerPoint COM.",
         ACCENT_CYAN),
    ]

    for num, title, desc, color in steps:
        y = Inches(1.3) + (num - 1) * Inches(1.15)
        add_numbered_step(slide, Inches(0.8), y, num, title, desc, accent_color=color)

        # Connector arrow
        if num < 5:
            arrow_y = y + Inches(0.85)
            add_text_box(slide, Inches(0.88), arrow_y, Inches(0.5), Inches(0.3),
                         "|", font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


def slide_research_deep_dive(prs):
    """Slide 5: Deep dive into the research module."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "Deep Dive: Research Module", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "research.py  -  928 lines  |  The heart of the application",
                 font_size=16, color=ACCENT_BLUE, bold=False, font_name="Consolas")

    # Left column: GeminiResearcher class
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.4),
                 "GeminiResearcher Class", font_size=18, color=ACCENT_CYAN, bold=True)

    methods = [
        "init_browser()  -  Launch Chrome with Selenium WebDriver",
        "navigate_to_gemini()  -  Go to gemini.google.com",
        "check_login_required()  -  Detect login state on Win11",
        "perform_deep_research()  -  Submit prompt & wait for AI",
        "wait_for_research_completion()  -  Smart polling loop",
        "extract_research_data()  -  Parse HTML response",
        "parse_research_text()  -  Structure into sections",
        "ask_gemini()  -  Send follow-up questions",
    ]

    add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.8), Inches(4.5),
                    methods, font_size=13, color=TEXT_GRAY, bullet_color=ACCENT_BLUE)

    # Right column: Research levels
    add_text_box(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(0.4),
                 "3 Research Depth Levels", font_size=18, color=ACCENT_PURPLE, bold=True)

    levels_card = add_shape_rect(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(3.8), BG_CARD)

    level_data = [
        ("Level 1: Basic", "Simple, educational content for\ngeneral audiences. Quick research.", ACCENT_GREEN),
        ("Level 2: Professional", "Technical, data-driven content\nwith 2024-2026 studies.", ACCENT_ORANGE),
        ("Level 3: Master Thesis", "Maximum depth, 5000+ words.\nAcademic-grade research.", ACCENT_BLUE),
    ]

    for i, (name, desc, color) in enumerate(level_data):
        y = Inches(2.8) + i * Inches(1.2)
        add_accent_bar(slide, Inches(7.2), y, Inches(0.05), Inches(0.8), color)
        add_text_box(slide, Inches(7.5), y, Inches(4.5), Inches(0.3),
                     name, font_size=15, color=color, bold=True)
        add_text_box(slide, Inches(7.5), y + Inches(0.35), Inches(4.5), Inches(0.5),
                     desc, font_size=12, color=TEXT_GRAY)


def slide_content_generation(prs):
    """Slide 6: Content generation module."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "Content Generation & Structuring", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "content_generator.py  -  Turns raw research into presentation-ready slides",
                 font_size=16, color=ACCENT_GREEN, font_name="Consolas")

    # ContentGenerator class methods
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.4),
                 "ContentGenerator Class", font_size=20, color=ACCENT_CYAN, bold=True)

    steps_content = [
        "create_outline()  -  Generate structure from research",
        "create_outline_from_research()  -  Parse sections into slides",
        "ask_gemini_for_outline()  -  Fallback: ask Gemini directly",
        "parse_outline()  -  Convert text to structured slide data",
        "enhance_with_speaker_notes()  -  Add notes to each slide",
        "suggest_visualizations()  -  Ask Gemini for chart ideas",
        "parse_viz_suggestions()  -  Structure visualization data",
    ]

    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.8), Inches(3.5),
                    steps_content, font_size=14, color=TEXT_GRAY, bullet_color=ACCENT_GREEN)

    # Right side: Slide structure
    add_text_box(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(0.4),
                 "Generated Presentation Structure", font_size=18, color=ACCENT_PURPLE, bold=True)

    structure = [
        "Title Slide  -  Topic name + subtitle",
        "Table of Contents  -  Section overview",
        "Introduction  -  2-3 context slides",
        "Main Content  -  15-20 organized slides",
        "Data Visualizations  -  Charts & images",
        "Conclusion  -  Summary & takeaways",
        "References  -  Research citations",
        "Thank You  -  Closing slide",
    ]

    add_bullet_list(slide, Inches(7), Inches(2.7), Inches(5.8), Inches(4),
                    structure, font_size=14, color=TEXT_GRAY, bullet_color=ACCENT_PURPLE)


def slide_ppt_builder(prs):
    """Slide 7: PPT builder and visualizations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "PowerPoint Assembly & Visualizations", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_ORANGE)

    # Left: PPT Builder
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "ppt_builder.py", font_size=20, color=ACCENT_PURPLE, bold=True, font_name="Consolas")

    builder_items = [
        "add_title_slide()  -  Professional title with styling",
        "add_content_slide()  -  Smart overflow handling",
        "add_image_slide()  -  Proper image scaling",
        "add_two_column_slide()  -  Side-by-side layout",
        "build_presentation()  -  Orchestrates all slides",
        "Max 5 bullets/slide, auto-creates continuations",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3),
                    builder_items, font_size=14, color=TEXT_GRAY, bullet_color=ACCENT_PURPLE)

    # Right: Visualizations
    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                 "visualizations.py", font_size=20, color=ACCENT_ORANGE, bold=True, font_name="Consolas")

    viz_items = [
        "Google Images search via Selenium",
        "Downloads relevant diagrams & images",
        "Fallback: matplotlib sample charts",
        "Bar, Pie, Line chart generation",
        "PNG export at high resolution",
        "Professional color scheme matching",
    ]
    add_bullet_list(slide, Inches(7), Inches(2.0), Inches(5.8), Inches(3),
                    viz_items, font_size=14, color=TEXT_GRAY, bullet_color=ACCENT_ORANGE)

    # Bottom: Themes
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4),
                 "10 Built-in Themes", font_size=18, color=ACCENT_CYAN, bold=True)

    themes = ["Office Default", "Ion", "Integral", "Facet", "Wisp",
              "Retrospect", "Slice", "Organic", "Frame", "Basis"]

    for i, theme in enumerate(themes):
        x = Inches(0.8) + i * Inches(1.2)
        colors = [ACCENT_BLUE, ACCENT_PURPLE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_ORANGE,
                  ACCENT_BLUE, ACCENT_PURPLE, ACCENT_GREEN, ACCENT_CYAN, ACCENT_ORANGE]
        badge = add_shape_rect(slide, x, Inches(5.7), Inches(1.1), Inches(0.9), BG_CARD)
        add_text_box(slide, x + Inches(0.05), Inches(5.75), Inches(1), Inches(0.35),
                     str(i + 1), font_size=22, color=colors[i], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(6.1), Inches(1), Inches(0.4),
                     theme, font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)


def slide_super_research(prs):
    """Slide 8: Super Research & Scientific APIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "Super Research Mode", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_CYAN)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
                 "When enabled, augments Gemini research with data from academic APIs.\n"
                 "Also generates a supplemental Word document with extra research references.",
                 font_size=16, color=TEXT_GRAY)

    # API Cards
    apis = [
        ("CORE API", "Free academic paper\nsearch engine", "No API key needed\n3+ papers per topic", ACCENT_BLUE),
        ("arXiv", "Preprint server for\nscience & math", "Searches abstracts\nand metadata", ACCENT_PURPLE),
        ("PubMed", "Biomedical & life\nscience literature", "Medical research\npapers", ACCENT_GREEN),
        ("OpenLibrary", "Book metadata &\nISBN lookup", "Full book info\nand covers", ACCENT_ORANGE),
        ("Internet Archive", "Full text downloads\nof public books", "Free book content\nextraction", ACCENT_CYAN),
        ("LibGen", "PDF downloads of\nacademic books", "Fallback source\nfor textbooks", ACCENT_BLUE),
    ]

    for i, (name, desc, detail, color) in enumerate(apis):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4)
        y = Inches(2.5) + row * Inches(2.2)

        card = add_shape_rect(slide, x, y, Inches(3.6), Inches(1.9), BG_CARD)
        add_accent_bar(slide, x, y, Inches(3.6), Inches(0.05), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.2), Inches(0.35),
                     name, font_size=18, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.2), Inches(0.5),
                     desc, font_size=13, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.2), Inches(0.5),
                     detail, font_size=11, color=TEXT_DIM)


def slide_gui_wizard(prs):
    """Slide 9: GUI Wizard walkthrough."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "GUI Wizard  -  ppt_wizard.py", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Professional 6-step Tkinter wizard (1106 lines)  -  850 x 750px window",
                 font_size=16, color=ACCENT_PURPLE, font_name="Consolas")

    wizard_steps = [
        (0, "Welcome Screen", "Branding, version info, and 'Get Started' button", ACCENT_BLUE),
        (1, "Topic Input", "Enter topic + Gemini login check with retry", ACCENT_PURPLE),
        (2, "Advanced Options", "Slide count (10-30) & research depth (3 levels)", ACCENT_GREEN),
        (3, "Super Research", "Toggle super research & auto-images features", ACCENT_ORANGE),
        (4, "Theme Selection", "Visual gallery of 10 themes with scrolling", ACCENT_CYAN),
        (5, "Export Options", "Choose PPTX, PPT, or PDF export format", ACCENT_BLUE),
        (6, "Progress & Completion", "Real-time log, animated progress bar, open folder", ACCENT_PURPLE),
    ]

    for step_num, title, desc, color in wizard_steps:
        y = Inches(2.1) + step_num * Inches(0.72)
        add_numbered_step(slide, Inches(0.8), y, step_num, title, desc,
                          accent_color=color, width=Inches(10))


def slide_tech_stack(prs):
    """Slide 10: Technology stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "Technology Stack", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_ORANGE)

    categories = [
        ("Core Language", [
            ("Python 3.8+", "Main programming language"),
        ], ACCENT_BLUE),
        ("Browser Automation", [
            ("Selenium WebDriver", "Chrome browser control"),
            ("webdriver-manager", "Auto ChromeDriver setup"),
            ("BeautifulSoup4", "HTML parsing & extraction"),
        ], ACCENT_PURPLE),
        ("Presentation", [
            ("python-pptx", "PowerPoint PPTX creation"),
            ("comtypes", "COM automation for PPT/PDF"),
        ], ACCENT_GREEN),
        ("Data & Visuals", [
            ("matplotlib", "Chart & graph generation"),
            ("Pillow (PIL)", "Image processing"),
            ("requests", "HTTP API calls"),
        ], ACCENT_ORANGE),
        ("User Interface", [
            ("tkinter", "GUI wizard framework"),
            ("colorama", "Colored terminal output"),
            ("tqdm", "Progress bar display"),
        ], ACCENT_CYAN),
    ]

    for i, (category, items, color) in enumerate(categories):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4)
        y = Inches(1.5) + row * Inches(2.8)

        card = add_shape_rect(slide, x, y, Inches(3.6), Inches(2.4), BG_CARD)
        add_accent_bar(slide, x, y, Inches(3.6), Inches(0.05), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(0.35),
                     category, font_size=18, color=color, bold=True)

        for j, (lib, desc) in enumerate(items):
            ly = y + Inches(0.6) + j * Inches(0.55)
            add_text_box(slide, x + Inches(0.3), ly, Inches(3), Inches(0.25),
                         lib, font_size=14, color=TEXT_WHITE, bold=True, font_name="Consolas")
            add_text_box(slide, x + Inches(0.3), ly + Inches(0.25), Inches(3), Inches(0.25),
                         desc, font_size=11, color=TEXT_DIM)


def slide_user_flow(prs):
    """Slide 11: User flow / how to run."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                 "User Flow  -  From Topic to Presentation", font_size=36, color=TEXT_WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ACCENT_CYAN)

    flow_steps = [
        ("1.  Launch the app", "Run ppt_generator.py (CLI) or ppt_wizard.py (GUI), "
         "or use the packaged PPT_Generator.exe", ACCENT_BLUE),
        ("2.  Enter your topic", "Type any topic  -  e.g. 'Quantum Computing', "
         "'Climate Change', 'Digital Marketing'", ACCENT_PURPLE),
        ("3.  Configure options", "Choose slide count (10-30), research depth (Basic / Professional / "
         "Master), theme, and export format", ACCENT_GREEN),
        ("4.  Login to Gemini", "Chrome opens automatically. Log into your Google account "
         "if needed, then press Enter", ACCENT_ORANGE),
        ("5.  Wait 4-6 minutes", "Research: ~2-3 min  |  Content: ~1-2 min  |  "
         "Visuals: ~30s  |  Build: ~30s", ACCENT_CYAN),
        ("6.  Get your PPT!", "Presentation auto-saved to output/ folder. "
         "Folder opens automatically when done.", ACCENT_GREEN),
    ]

    for i, (title, desc, color) in enumerate(flow_steps):
        y = Inches(1.4) + i * Inches(0.95)
        card = add_shape_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), BG_CARD)
        add_accent_bar(slide, Inches(0.8), y, Inches(0.06), Inches(0.8), color)
        add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(10.5), Inches(0.35),
                     title, font_size=18, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.4), Inches(10.5), Inches(0.35),
                     desc, font_size=13, color=TEXT_GRAY)


def slide_closing(prs):
    """Slide 12: Thank you / closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(0.5),
                 "NOTEBOOKLM SOURCE GUIDE", font_size=14, color=ACCENT_CYAN, bold=True)

    add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(1),
                 "Thank You!", font_size=52, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(3.9), Inches(11), Inches(0.8),
                 "This presentation was generated to explain how the AI\n"
                 "PowerPoint Generator works  -  from research to final slides.",
                 font_size=20, color=TEXT_GRAY)

    add_accent_bar(slide, Inches(1), Inches(5.0), Inches(2), Inches(0.04), ACCENT_BLUE)

    details = [
        "Made by Luqman  (2024-ag-8738 DVM)",
        "Python  +  Selenium  +  Google Gemini  +  python-pptx",
        "Total codebase:  ~4,000 lines across 10+ files",
        "Discord:  sad_memer.",
    ]
    add_bullet_list(slide, Inches(1), Inches(5.3), Inches(10), Inches(2),
                    details, font_size=15, color=TEXT_GRAY, bullet_color=ACCENT_PURPLE)

    # Bottom accent
    add_accent_bar(slide, Inches(0), Inches(7.3), SLIDE_WIDTH, Inches(0.2), ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    slide_title(prs)           # 1. Title
    slide_overview(prs)        # 2. What it does
    slide_architecture(prs)    # 3. Architecture
    slide_pipeline(prs)        # 4. 5-step pipeline
    slide_research_deep_dive(prs)  # 5. Research module
    slide_content_generation(prs)  # 6. Content generation
    slide_ppt_builder(prs)     # 7. PPT builder & visuals
    slide_super_research(prs)  # 8. Super research
    slide_gui_wizard(prs)      # 9. GUI wizard
    slide_tech_stack(prs)      # 10. Tech stack
    slide_user_flow(prs)       # 11. User flow
    slide_closing(prs)         # 12. Thank you

    # Save
    output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "NotebookLM_How_PPT_Generator_Works.pptx")
    prs.save(output_path)
    print(f"\n  Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Format: 16:9 widescreen\n")


if __name__ == "__main__":
    main()
